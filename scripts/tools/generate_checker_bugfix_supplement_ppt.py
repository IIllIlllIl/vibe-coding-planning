"""Generate PPT slides for the Bug Fix subset of the original 155-case cohort."""

from __future__ import annotations

import datetime
import hashlib
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from scripts.tools.generate_checker_interim_ppt import (  # noqa: E402
    ARMS,
    NAVY,
    WHITE,
    add_bullets,
    add_chart,
    add_footer,
    add_metric_table,
    add_title,
    metrics,
    style_table,
)

ROOT = Path("output/checker_eval/polybench-flash-pro-kimi-baseline")
BUG_FIX_PATH = Path(
    "output/SWE-PolyBench/polybench-pct-checker-datasets/"
    "20260609_198_cdf4d414e401/derived/task_category_v1/"
    "bug_fix_cases.jsonl"
)
OUTPUT = Path("checker_only_bugfix_subset_supplement_2026-06-11.pptx")
ORIGINAL_SNAPSHOT_CUTOFF = datetime.datetime(
    2026, 6, 11, 9, 52, 0
).timestamp()
EXPECTED_COMMON_COHORT = 155


def _load_original_snapshot_predictions(
    arm: str,
) -> dict[str, dict[str, Any]]:
    predictions: dict[str, dict[str, Any]] = {}
    for path in sorted((ROOT / arm / "instances").glob("*/prediction.json")):
        if path.stat().st_mtime > ORIGINAL_SNAPSHOT_CUTOFF:
            continue
        row = json.loads(path.read_text(encoding="utf-8"))
        predictions[row["instance_id"]] = row
    return predictions


def _load_bug_fix_ids() -> set[str]:
    return {
        json.loads(line)["instance_id"]
        for line in BUG_FIX_PATH.read_text(encoding="utf-8").splitlines()
        if line.strip()
    }


def _cohort() -> tuple[
    dict[str, dict[str, dict[str, Any]]],
    set[str],
    set[str],
]:
    rows = {
        arm: _load_original_snapshot_predictions(arm) for arm, _ in ARMS
    }
    common = set.intersection(*(set(rows[arm]) for arm, _ in ARMS))
    if len(common) != EXPECTED_COMMON_COHORT:
        raise ValueError(
            f"Expected original common cohort of 155, found {len(common)}"
        )
    bug_fix = _load_bug_fix_ids()
    return rows, common, common & bug_fix


def _add_confusion_table(
    slide: Any, data: dict[str, dict[str, float | int]]
) -> None:
    headers = ["Arm", "TP", "FP", "FN", "TN", "Recall", "Specificity"]
    shape = slide.shapes.add_table(
        5, 7, Inches(1.0), Inches(1.45), Inches(11.3), Inches(2.8)
    )
    table = shape.table
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, (arm, label) in enumerate(ARMS, 1):
        result = data[arm]
        values = [
            label,
            str(result["tp"]),
            str(result["fp"]),
            str(result["fn"]),
            str(result["tn"]),
            f"{result['recall']:.1%}",
            f"{result['specificity']:.1%}",
        ]
        for col, value in enumerate(values):
            table.cell(row_index, col).text = value
    style_table(table)


def main() -> None:
    rows, common_ids, bug_fix_ids = _cohort()
    common_metrics = {
        arm: metrics(rows[arm][instance] for instance in common_ids)
        for arm, _ in ARMS
    }
    bug_fix_metrics = {
        arm: metrics(rows[arm][instance] for instance in bug_fix_ids)
        for arm, _ in ARMS
    }
    resolved = sum(
        rows["flash_rules"][instance]["test_results"]["resolved"]
        for instance in bug_fix_ids
    )
    cohort_sha = hashlib.sha256(
        ("\n".join(sorted(bug_fix_ids)) + "\n").encode("utf-8")
    ).hexdigest()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.3)
    )
    p = title.text_frame.paragraphs[0]
    p.text = "Bug Fix Subset Analysis"
    p.font.name = "Aptos Display"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    subtitle = slide.shapes.add_textbox(
        Inches(0.85), Inches(3.4), Inches(11.6), Inches(1.4)
    )
    p = subtitle.text_frame.paragraphs[0]
    p.text = (
        "Official PolyBench task_category label\n"
        "Intersection with the original 155-case common cohort"
    )
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(210, 222, 235)

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Cohort Definition",
        "This supplement preserves the exact sample boundary used in the previous deck",
    )
    add_bullets(
        slide,
        [
            "Original common cohort: 155 cases completed successfully by all four arms.",
            "Official PolyBench Bug Fix cases in the 198-case snapshot: 156.",
            f"Intersection used here: {len(bug_fix_ids)} Bug Fix cases.",
            f"Ground truth distribution: {resolved} resolved and {len(bug_fix_ids) - resolved} unresolved.",
            "Recovery outputs written after the previous PPT snapshot are intentionally excluded.",
            f"Reproducibility cohort SHA256: {cohort_sha[:16]}…",
        ],
        size=18,
    )
    add_footer(
        slide,
        "Official label source: AmazonScience/SWE-PolyBench.task_category",
    )

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Bug Fix Subset Metrics",
        f"Paired comparison on the same {len(bug_fix_ids)} cases",
    )
    add_metric_table(slide, bug_fix_metrics, top=1.25)
    add_chart(
        slide,
        [label for _, label in ARMS],
        [
            (
                "Balanced accuracy",
                [
                    float(bug_fix_metrics[arm]["balanced_accuracy"])
                    for arm, _ in ARMS
                ],
            ),
            (
                "Recall",
                [float(bug_fix_metrics[arm]["recall"]) for arm, _ in ARMS],
            ),
            (
                "F1",
                [float(bug_fix_metrics[arm]["f1"]) for arm, _ in ARMS],
            ),
        ],
        left=1.05,
        top=4.05,
        width=11.2,
        height=2.6,
    )
    add_footer(slide, "Bug Fix ∩ original common cohort; N = 125")

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Bug Fix vs. Full Common Cohort",
        "Recall and balanced accuracy before and after restricting task type",
    )
    add_chart(
        slide,
        [label for _, label in ARMS],
        [
            (
                "Full recall (N=155)",
                [
                    float(common_metrics[arm]["recall"])
                    for arm, _ in ARMS
                ],
            ),
            (
                "Bug Fix recall (N=125)",
                [
                    float(bug_fix_metrics[arm]["recall"])
                    for arm, _ in ARMS
                ],
            ),
            (
                "Full balanced accuracy",
                [
                    float(common_metrics[arm]["balanced_accuracy"])
                    for arm, _ in ARMS
                ],
            ),
            (
                "Bug Fix balanced accuracy",
                [
                    float(bug_fix_metrics[arm]["balanced_accuracy"])
                    for arm, _ in ARMS
                ],
            ),
        ],
        left=0.8,
        top=1.45,
        width=7.8,
        height=4.8,
    )
    add_bullets(
        slide,
        [
            "Recall remains highest for the no-rule baseline.",
            "Flash retains similar recall and F1 on Bug Fix tasks.",
            "Pro remains highly conservative with 3.6% recall.",
            "No arm exceeds 0.48 balanced accuracy.",
        ],
        left=8.85,
        top=1.7,
        width=3.8,
        height=4.2,
        size=14,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Bug Fix Confusion Matrix")
    _add_confusion_table(slide, bug_fix_metrics)
    add_bullets(
        slide,
        [
            "Flash and no-rule baseline remain permissive, producing many false positives.",
            "Pro remains strongly conservative: 2 true positives and 53 false negatives.",
            "All MCC values remain negative, indicating no positive prediction correlation.",
        ],
        top=4.65,
        height=1.8,
        size=15,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "Restricting the analysis to officially labeled Bug Fix tasks does not improve discrimination.",
            "The mixed-task composition of PolyBench is therefore unlikely to be the primary explanation for the near-random aggregate result.",
            "The main observed effect is still checker strictness: Flash/no-rules are permissive, while Pro is highly conservative.",
            "Further analysis should examine rule applicability, individual rule triggers, repository effects, and checker calibration.",
        ],
        size=19,
    )
    add_footer(
        slide,
        "Interim analysis; final complete-cohort metrics should be recomputed after recovery",
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
