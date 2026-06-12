"""Generate an English PPT for the interim checker-only results."""

from __future__ import annotations

import json
import math
from datetime import datetime
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("output/checker_eval/polybench-flash-pro-kimi-baseline")
OUTPUT = Path("checker_only_interim_results_2026-06-11.pptx")
ARMS = (
    ("flash_rules", "Flash rules"),
    ("no_rules", "No rules"),
    ("pro_rules", "Pro rules"),
    ("kimi_rules", "Kimi rules"),
)

NAVY = RGBColor(24, 45, 74)
BLUE = RGBColor(49, 105, 171)
ORANGE = RGBColor(242, 142, 43)
GREEN = RGBColor(89, 161, 79)
RED = RGBColor(225, 87, 89)
GRAY = RGBColor(105, 117, 130)
LIGHT = RGBColor(239, 243, 247)
WHITE = RGBColor(255, 255, 255)
COLORS = (BLUE, GRAY, ORANGE, GREEN)


def load_predictions(arm: str) -> dict[str, dict[str, Any]]:
    paths = sorted((ROOT / arm / "instances").glob("*/prediction.json"))
    return {
        item["instance_id"]: item
        for item in (
            json.loads(path.read_text(encoding="utf-8")) for path in paths
        )
    }


def metrics(rows: Iterable[dict[str, Any]]) -> dict[str, float | int]:
    tp = fp = fn = tn = 0
    for row in rows:
        predicted = row["check_result"]["passed"]
        resolved = row["test_results"]["resolved"]
        if predicted and resolved:
            tp += 1
        elif predicted:
            fp += 1
        elif resolved:
            fn += 1
        else:
            tn += 1
    total = tp + fp + fn + tn
    precision = tp / (tp + fp) if tp + fp else 0.0
    recall = tp / (tp + fn) if tp + fn else 0.0
    specificity = tn / (tn + fp) if tn + fp else 0.0
    f1 = (
        2 * precision * recall / (precision + recall)
        if precision + recall
        else 0.0
    )
    denominator = math.sqrt(
        (tp + fp) * (tp + fn) * (tn + fp) * (tn + fn)
    )
    return {
        "n": total,
        "tp": tp,
        "fp": fp,
        "fn": fn,
        "tn": tn,
        "accuracy": (tp + tn) / total if total else 0.0,
        "precision": precision,
        "recall": recall,
        "f1": f1,
        "specificity": specificity,
        "balanced_accuracy": (recall + specificity) / 2,
        "mcc": (
            (tp * tn - fp * fn) / denominator if denominator else 0.0
        ),
        "pass_rate": (tp + fp) / total if total else 0.0,
    }


def add_title(slide: Any, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.65)
    )
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.58), Inches(0.9), Inches(12), Inches(0.4)
        )
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY


def add_footer(slide: Any, text: str = "Interim snapshot · incomplete/error cases excluded") -> None:
    line = slide.shapes.add_shape(
        1, Inches(0.55), Inches(7.18), Inches(12.2), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()
    box = slide.shapes.add_textbox(
        Inches(0.58), Inches(7.22), Inches(12), Inches(0.2)
    )
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY


def add_bullets(
    slide: Any,
    items: list[str],
    *,
    left: float = 0.8,
    top: float = 1.45,
    width: float = 11.7,
    height: float = 5.3,
    size: int = 20,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(13)
        p.text = f"• {item}"


def style_table(table: Any, header_color: RGBColor = NAVY) -> None:
    for col in range(len(table.columns)):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
    for row in range(1, len(table.rows)):
        for col in range(len(table.columns)):
            cell = table.cell(row, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row % 2 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(10)
                p.font.color.rgb = NAVY
                p.alignment = PP_ALIGN.CENTER


def add_metric_table(
    slide: Any,
    data: dict[str, dict[str, float | int]],
    *,
    top: float,
) -> None:
    headers = [
        "Arm",
        "N",
        "Accuracy",
        "Precision",
        "Recall",
        "F1",
        "Bal. Acc.",
        "MCC",
        "Pass rate",
    ]
    shape = slide.shapes.add_table(
        len(ARMS) + 1,
        len(headers),
        Inches(0.55),
        Inches(top),
        Inches(12.2),
        Inches(2.55),
    )
    table = shape.table
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, (arm, label) in enumerate(ARMS, 1):
        result = data[arm]
        values = [
            label,
            str(result["n"]),
            f"{result['accuracy']:.3f}",
            f"{result['precision']:.3f}",
            f"{result['recall']:.3f}",
            f"{result['f1']:.3f}",
            f"{result['balanced_accuracy']:.3f}",
            f"{result['mcc']:.3f}",
            f"{result['pass_rate']:.3f}",
        ]
        for col, value in enumerate(values):
            table.cell(row_index, col).text = value
    table.columns[0].width = Inches(1.55)
    style_table(table)


def add_chart(
    slide: Any,
    categories: list[str],
    series: list[tuple[str, list[float]]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    maximum: float = 1.0,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = maximum
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.has_title = False
    for index, chart_series in enumerate(chart.series):
        chart_series.format.fill.solid()
        chart_series.format.fill.fore_color.rgb = COLORS[index % len(COLORS)]


def main() -> None:
    rows = {arm: load_predictions(arm) for arm, _ in ARMS}
    common_ids = set.intersection(*(set(rows[arm]) for arm, _ in ARMS))
    available = {
        arm: metrics(rows[arm].values()) for arm, _ in ARMS
    }
    common = {
        arm: metrics(rows[arm][instance] for instance in common_ids)
        for arm, _ in ARMS
    }
    snapshot = datetime.now().astimezone().strftime("%Y-%m-%d %H:%M %Z")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    banner = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.15), Inches(11.7), Inches(1.35)
    )
    p = title.text_frame.paragraphs[0]
    p.text = "Checker-Only Rule Evaluation"
    p.font.name = "Aptos Display"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(
        Inches(0.85), Inches(3.55), Inches(11.6), Inches(1.1)
    )
    p = sub.text_frame.paragraphs[0]
    p.text = (
        "Interim PolyBench results on completed, successful checker runs\n"
        f"Snapshot: {snapshot}"
    )
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(210, 222, 235)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Evaluation Design")
    add_bullets(
        slide,
        [
            "Input: 198 fixed PCT plans with existing resolved/unresolved labels.",
            "Checker model: DeepSeek Flash for all four arms.",
            "Arms: Flash rules, no-rule baseline, Pro rules, and Kimi rules.",
            "Positive prediction: the checker accepts the plan; ground truth positive: PCT resolved.",
            "This snapshot excludes checker errors and cases not yet completed.",
            "Two views are reported: all available predictions per arm and a common 155-case cohort.",
        ],
        size=18,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Snapshot Coverage",
        "Successful predictions available at snapshot time; excluded cases are not scored",
    )
    counts = [int(available[arm]["n"]) for arm, _ in ARMS]
    add_chart(
        slide,
        [label for _, label in ARMS],
        [("Completed", counts)],
        left=0.7,
        top=1.4,
        width=7.4,
        height=4.8,
        maximum=210,
    )
    missing = [198 - count for count in counts]
    add_bullets(
        slide,
        [
            f"Flash rules: {counts[0]}/198 complete",
            f"No rules: {counts[1]}/198 complete; {missing[1]} excluded",
            f"Pro rules: {counts[2]}/198 complete; {missing[2]} excluded",
            f"Kimi rules: {counts[3]}/198 complete; {missing[3]} excluded",
            f"Common completed cohort: {len(common_ids)} cases",
        ],
        left=8.35,
        top=1.65,
        width=4.3,
        height=4.5,
        size=16,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "All Available Successful Runs",
        "Each arm uses its own currently completed sample set; N therefore differs",
    )
    add_metric_table(slide, available, top=1.35)
    add_bullets(
        slide,
        [
            "Pro rules have the highest raw accuracy, but accept only 9.3% of plans.",
            "Flash rules provide the strongest rule-based recall and F1 balance.",
            "No rules has the highest recall, paired with the highest false-positive tendency.",
            "Negative MCC values indicate weak discrimination in the current snapshot.",
        ],
        top=4.25,
        height=2.3,
        size=15,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Fair Comparison: Common 155-Case Cohort",
        "All four arms produced successful predictions for the same cases",
    )
    add_metric_table(slide, common, top=1.25)
    add_chart(
        slide,
        [label for _, label in ARMS],
        [
            ("Accuracy", [float(common[a]["accuracy"]) for a, _ in ARMS]),
            ("Recall", [float(common[a]["recall"]) for a, _ in ARMS]),
            ("F1", [float(common[a]["f1"]) for a, _ in ARMS]),
        ],
        left=1.05,
        top=4.05,
        width=11.2,
        height=2.6,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Confusion Matrix Counts: Common Cohort")
    headers = ["Arm", "TP", "FP", "FN", "TN", "Resolved recall", "Specificity"]
    shape = slide.shapes.add_table(
        5, 7, Inches(1.05), Inches(1.55), Inches(11.2), Inches(3.2)
    )
    table = shape.table
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, (arm, label) in enumerate(ARMS, 1):
        result = common[arm]
        values = [
            label,
            str(result["tp"]),
            str(result["fp"]),
            str(result["fn"]),
            str(result["tn"]),
            f"{result['recall']:.1%}",
            f"{result['specificity']:.1%}",
        ]
        for col, value in enumerate(values):
            table.cell(row_index, col).text = value
    style_table(table)
    add_bullets(
        slide,
        [
            "Pro rules reject most plans: only 3 true positives but 82 true negatives.",
            "Flash rules recover 40 of 61 resolved cases, while producing 65 false positives.",
            "The baseline accepts more plans, increasing recall but also false positives.",
        ],
        top=5.05,
        height=1.5,
        size=14,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Interpretation")
    add_bullets(
        slide,
        [
            "Accuracy alone is misleading under the current class distribution and checker behavior.",
            "Pro rules appear conservative rather than discriminative: high specificity, very low recall.",
            "Flash rules are more balanced than Pro and Kimi, but still over-accept unresolved plans.",
            "No-rule performance shows that the base agent already captures some plan-quality signal.",
            "Rule quality should be judged using F1, balanced accuracy, MCC, and paired case analysis.",
        ],
        size=19,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Limitations and Next Update")
    add_bullets(
        slide,
        [
            "This is an interim snapshot, not the final four-arm comparison.",
            "Excluded cases: 17 no-rule, 4 Pro, and 24 Kimi predictions at snapshot time.",
            "Infrastructure failures and incomplete runs are omitted rather than counted as checker failures.",
            "The active recovery run will fill missing cases without rerunning successful predictions.",
            "The final report will recompute all metrics on the complete 198-case paired cohort.",
        ],
        size=18,
    )
    add_footer(
        slide,
        "Source: output/checker_eval/polybench-flash-pro-kimi-baseline",
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
